from pptx import Presentation
import os

def remover_arquivo(path: str):
    if os.path.exists(path):
        os.remove(path)


def substituir_no_paragrafo(paragraph, dados):

    texto_paragrafo = "".join([run.text for run in paragraph.runs])

    modificado = False

    for chave, valor in dados.items():

        marcador = f"{{{{{chave}}}}}"

        if marcador in texto_paragrafo:
            texto_paragrafo = texto_paragrafo.replace(marcador, str(valor))
            modificado = True

    if modificado and len(paragraph.runs) > 0:

        paragraph.runs[0].text = texto_paragrafo

        for i in range(1, len(paragraph.runs)):
            paragraph.runs[i].text = ""


def preencher_pptx(template_path: str, output_path: str, dados: dict):

    prs = Presentation(template_path)

    for slide in prs.slides:

        for shape in slide.shapes:

            if shape.has_text_frame:

                for paragraph in shape.text_frame.paragraphs:
                    substituir_no_paragrafo(paragraph, dados)

            if shape.has_table:

                for row in shape.table.rows:

                    for cell in row.cells:

                        for paragraph in cell.text_frame.paragraphs:
                            substituir_no_paragrafo(paragraph, dados)

    prs.save(output_path)